# -*- coding: utf-8 -*-
"""
ÉHON 提案書ビルダー
基準: Gene_Story_Platform_—_総合事業計画.pdf のトンマナ

使い方:
    import sys; sys.path.insert(0, ".claude/lib")
    from ehon_deck import *
    prs = new_deck()
    s = cover(prs, eyebrow="ÉHON INC. × 〇〇株式会社", title="...", sub="...", date="2026年8月27日")
    s = slide(prs, "課題と機会")
    card(s, 0.6, 1.6, 5.8, 3.2, accent=NAVY, title="① ...", lines=["...","..."])
    conclusion(s, "言い切りの一文")
    prs.save("output/xxx.pptx")

守ること:
  * 色は下のパレット以外を使わない
  * 罫・矢印は文字で代用しない（この関数群が図形で描く）
  * イラストは入れない
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import copy

# ============ パレット（実測値・この外を使わない） ============
CREAM  = RGBColor(0xF4, 0xEF, 0xE7)   # 地色
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)   # カード地
WINE   = RGBColor(0x65, 0x32, 0x45)   # 主アクセント
NAVY   = RGBColor(0x21, 0x2A, 0x42)   # 副アクセント
ROSE   = RGBColor(0x90, 0x4E, 0x65)   # 中間調
GOLD   = RGBColor(0xC4, 0xAA, 0x6B)   # 罫・矢印・第3カテゴリ
GOLD_L = RGBColor(0xDC, 0xCA, 0x9B)   # 淡金
RULE   = RGBColor(0xE0, 0xDC, 0xD7)   # 罫
MUTED  = RGBColor(0x5C, 0x5B, 0x58)   # 注記
INK    = RGBColor(0x21, 0x2A, 0x42)   # 本文（濃紺を墨として使う）

# ============ 書体 ============
MINCHO = "ヒラギノ明朝 ProN W6"        # 見出し
MINCHO_M = "ヒラギノ明朝 ProN W3"
GOTHIC = "ヒラギノ角ゴシック W3"        # 本文
GOTHIC_B = "ヒラギノ角ゴシック W6"

W, H = 13.333, 7.5                     # 16:9
FOOTER = "Confidential ／ ÉHON INC."


def _set_font(run, name, size, color, bold=False):
    """latin と eastasia の両方に書体を当てる（片方だけだと日本語が既定書体に落ちる）"""
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = name
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:ea", "a:cs"):
        e = rPr.find(qn(tag))
        if e is None:
            e = rPr.makeelement(qn(tag), {})
            rPr.append(e)
        e.set("typeface", name)



def _est_lines(text, width_in, pt):
    """全角=pt/72インチ、半角=その半分として折り返し行数を見積もる"""
    if not text:
        return 1
    em = pt / 72.0
    width = 0.0
    for ch in str(text):
        # 全角は em。半角でも数字・英字は 0.55em 程度あるので、0.5 だと足りない
        width += em if ord(ch) > 0x2000 else em * 0.56
    import math
    # 実測より 1 行少なく見積もると文字が重なる。安全側に 6% 詰めて数える
    usable = max(width_in, 0.1) * 0.94
    return max(1, math.ceil(width / usable - 1e-9))


def _rect(slide, x, y, w, h, fill, line=None, line_w=0.75):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.shadow.inherit = False
    if fill is None:
        sh.fill.background()
    else:
        sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line; sh.line.width = Pt(line_w)
    return sh


def _text(slide, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, spacing=1.5):
    """runs = [(文字列, 書体, 級数, 色, 太字), ...] または [[...同段落内...], ...]"""
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    paras = runs if runs and isinstance(runs[0], list) else [[r] for r in runs]
    for i, para in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = spacing
        for (txt, fnt, sz, col, bold) in para:
            r = p.add_run(); r.text = txt
            _set_font(r, fnt, sz, col, bold)
    return tb



# ============ 縦位置の自動送り ============
# 座標を手打ちすると要素が重なる（実際に注記ブロックがカードに覆いかぶさる事故が出た）。
# slide() がカーソルを持ち、y を省略した部品はその位置に積まれて自動で下へ送られる。

MARGIN_X = 0.62
BODY_W = W - MARGIN_X * 2
BOTTOM = H - 1.52          # 結論帯の上端。ここを越えたら詰め込みすぎ
GAP = 0.24


def _cursor(slide):
    return getattr(slide, "_cur_y", 1.55)


def _advance(slide, h, gap=GAP):
    slide._cur_y = _cursor(slide) + h + gap


def space(slide, h=0.2):
    """意図的に余白を入れる"""
    slide._cur_y = _cursor(slide) + h


def remaining(slide):
    """結論帯までに残っている高さ"""
    return BOTTOM - _cursor(slide)


def overflowing(slide):
    return _cursor(slide) > BOTTOM


def assert_fits(prs):
    """保存前に必ず呼ぶ。はみ出しているスライドがあれば例外にする。
    検出できるだけでは防げない（実際に結論帯を突き抜けた提案書が出た）。"""
    bad = []
    for i, s in enumerate(prs.slides, 1):
        if getattr(s, "_cur_y", None) is None:
            continue
        over = _cursor(s) - BOTTOM
        if over > 0.02:
            bad.append(f"  {i}枚目「{getattr(s, '_heading', '')}」が {over:.2f}インチはみ出し")
    if bad:
        raise RuntimeError(
            "スライドがはみ出しています。項目を減らすか、スライドを分けてください。\n"
            + "\n".join(bad))
    return True


# ============ スライド ============
def new_deck():
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(W), Inches(H)
    return prs


def _blank(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _rect(s, 0, 0, W, H, CREAM)
    return s


def footer(slide):
    _text(slide, W - 3.4, H - 0.5, 3.0, 0.3,
          [(FOOTER, GOTHIC, 9, MUTED, False)], align=PP_ALIGN.RIGHT)


def cover(prs, eyebrow, title, sub, date):
    """表紙：左にワインの帯、大きなセリフタイトル、金の短い罫"""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _rect(s, 0, 0, W, H, CREAM)
    _rect(s, 0, 0, 0.42, H, WINE)                       # 左の帯
    # タイトルは折り返し行数に応じて級数を落とし、後続要素を押し下げる
    tw = 11.0
    pt = 40
    while _est_lines(title, tw, pt) > 2 and pt > 28:
        pt -= 2
    t_lines = _est_lines(title, tw, pt)
    t_h = t_lines * (pt / 72.0) * 1.25

    y = 2.15
    _text(s, 1.25, y, 10.5, 0.4, [(eyebrow, GOTHIC_B, 15, WINE, True)])
    y += 0.50
    _text(s, 1.25, y, tw, t_h, [(title, MINCHO, pt, NAVY, True)], spacing=1.2)
    y += t_h + 0.22
    if sub:
        s_h = _est_lines(sub, tw, 21) * 0.36
        _text(s, 1.25, y, tw, s_h, [(sub, MINCHO_M, 21, NAVY, False)], spacing=1.3)
        y += s_h + 0.32
    _rect(s, 1.25, y, 1.05, 0.035, GOLD)                # 金の短い罫
    y += 0.42
    _text(s, 1.25, y, 8.0, 0.3, [(date, GOTHIC, 12, MUTED, False)])
    footer(s)
    return s


def slide(prs, heading, sub=None):
    """本文スライド：見出し左のワイン縦バー（この資料の顔）"""
    s = _blank(prs)
    s._heading = heading
    _rect(s, 0.62, 0.42, 0.07, 0.42, WINE)              # 縦バー
    _text(s, 0.86, 0.36, 11.6, 0.6, [(heading, MINCHO, 29, NAVY, True)], spacing=1.1)
    if sub:
        _text(s, 0.86, 1.02, 11.6, 0.35, [(sub, GOTHIC, 13, MUTED, False)])
        s._cur_y = 1.62
    else:
        s._cur_y = 1.38
    footer(s)
    return s


# ============ 部品 ============
def card(slide, x=None, y=None, w=None, h=None, title=None, lines=None, accent=NAVY, badge_text=None):
    """白カード＋上辺の色バー。h=None で中身に合わせて自動。
    タイトル・本文の折り返しを見積もって高さを決めるので、文字が重ならない。
    箇条書きの点は図形で描く（文字の■を使わない）。"""
    auto = y is None
    if x is None: x = MARGIN_X
    if w is None: w = BODY_W
    if auto: y = _cursor(slide)
    tw = w - 0.52 - (1.72 if badge_text else 0)
    t_lines = _est_lines(title, tw, 17) if title else 0
    title_h = t_lines * 0.30 + 0.16 if title else 0.0

    bw = w - 0.74
    l_lines = [_est_lines(ln, bw, 12) for ln in (lines or [])]
    l_heights = [n * 0.235 + 0.085 for n in l_lines]

    if h is None:
        h = 0.26 + title_h + sum(l_heights) + 0.24
    _rect(slide, x, y, w, h, WHITE)
    _rect(slide, x, y, w, 0.05, accent)                 # 上辺バー

    cy = y + 0.26
    if title:
        _text(slide, x + 0.26, cy, tw, title_h,
              [(title, MINCHO, 17, NAVY, True)], spacing=1.2)
        cy += title_h
    for k, ln in enumerate(lines or []):
        _rect(slide, x + 0.28, cy + 0.09, 0.075, 0.075, accent)   # 点は図形
        _text(slide, x + 0.46, cy, bw, l_heights[k],
              [(ln, GOTHIC, 12, INK, False)], spacing=1.15)
        cy += l_heights[k]
    if badge_text:
        badge(slide, x + w - 1.55, y + 0.26, 1.3, badge_text, accent)
    if auto: _advance(slide, h)
    return h



def card_height(w, title=None, lines=None, badge_text=None):
    """card() が必要とする高さを、描かずに見積もる"""
    tw = w - 0.52 - (1.72 if badge_text else 0)
    title_h = (_est_lines(title, tw, 17) * 0.30 + 0.16) if title else 0.0
    bw = w - 0.74
    lh = sum(_est_lines(ln, bw, 12) * 0.235 + 0.085 for ln in (lines or []))
    return 0.26 + title_h + lh + 0.24


def cards_row(slide, x=None, y=None, total_w=None, items=None, gap=0.2):
    """カードを横並びにし、高さを最も高いものに揃える。
    items = [{"title":..., "lines":[...], "accent":..., "badge_text":...}, ...]"""
    auto = y is None
    if x is None: x = MARGIN_X
    if total_w is None: total_w = BODY_W
    if auto: y = _cursor(slide)
    n = len(items)
    cw = (total_w - gap * (n - 1)) / n
    h = max(card_height(cw, it.get("title"), it.get("lines"), it.get("badge_text")) for it in items)
    for i, it in enumerate(items):
        card(slide, x + i * (cw + gap), y, cw, h,
             title=it.get("title"), lines=it.get("lines"),
             accent=it.get("accent", NAVY), badge_text=it.get("badge_text"))
    if auto: _advance(slide, h)
    return h


def badge(slide, x, y, w, text, color=NAVY):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.28))
    sh.shadow.inherit = False
    sh.fill.solid(); sh.fill.fore_color.rgb = color
    sh.line.fill.background()
    tf = sh.text_frame; tf.word_wrap = False
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text
    _set_font(r, GOTHIC_B, 9.5, WHITE, True)


def conclusion(slide, text):
    """最下部の全幅ワイン帯。1スライド1本まで"""
    _rect(slide, 0.62, H - 1.32, W - 1.24, 0.62, WINE)
    _text(slide, 0.62, H - 1.22, W - 1.24, 0.45,
          [(text, MINCHO, 18, WHITE, True)], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def big_number(slide, x, y, w, label, value, unit=None, note=None, color=WINE):
    """指標名（小）→ 数値（特大）→ 単位（小）。3〜4個の数字はグラフにせずこれ"""
    _text(slide, x, y, w, 0.3, [(label, GOTHIC, 12, MUTED, False)])
    runs = [(value, MINCHO, 46, color, True)]
    if unit:
        runs.append((" " + unit, GOTHIC, 13, MUTED, False))
    _text(slide, x, y + 0.34, w, 0.85, [runs], spacing=1.0)
    if note:
        _text(slide, x, y + 1.24, w, 0.5, [(note, GOTHIC, 11, INK, False)], spacing=1.45)


def split_bar(slide, x, y, w, parts, h=0.72):
    """1本の帯を2〜3色に割る。parts=[(ラベル, 値文字列, 比率, 色), ...]"""
    total = sum(p[2] for p in parts) or 1
    cx = x
    for (label, val, ratio, col) in parts:
        pw = w * ratio / total
        _rect(slide, cx, y, pw, h, col)
        _text(slide, cx + 0.16, y + 0.09, pw - 0.32, 0.22, [(label, GOTHIC, 10, WHITE, False)])
        _text(slide, cx + 0.16, y + 0.32, pw - 0.32, 0.34, [(val, MINCHO, 19, WHITE, True)], spacing=1.0)
        cx += pw


def note_block(slide, x=None, y=None, w=None, h=None, title="", body=""):
    """左辺だけワインの太罫。囲まない。y・h は省略で自動"""
    auto = y is None
    if x is None: x = MARGIN_X
    if w is None: w = BODY_W
    if auto: y = _cursor(slide)
    if h is None:
        h = 0.52 + _est_lines(body, w - 0.62, 11.5) * 0.26 + 0.24
    _rect(slide, x, y, w, h, GOLD_L if False else RGBColor(0xEE, 0xE7, 0xDD))
    _rect(slide, x, y, 0.045, h, WINE)
    _text(slide, x + 0.24, y + 0.16, w - 0.48, 0.3, [(title, GOTHIC_B, 12, WINE, True)])
    _text(slide, x + 0.28, y + 0.52, w - 0.62, h - 0.72, [(body, GOTHIC, 11.5, INK, False)], spacing=1.5)
    if auto: _advance(slide, h)
    return h


def flow(slide, x=None, y=None, w=None, items=None, h=0.62, gap=0.22):
    """横並びのフロー。間は金の三角でつなぐ（矢印を文字で描かない）"""
    auto = y is None
    if x is None: x = MARGIN_X
    if w is None: w = BODY_W
    if auto: y = _cursor(slide)
    n = len(items)
    bw = (w - gap * (n - 1)) / n
    for i, t in enumerate(items):
        bx = x + i * (bw + gap)
        _rect(slide, bx, y, bw, h, NAVY)
        _text(slide, bx + 0.08, y + 0.06, bw - 0.16, h - 0.12,
              [(t, GOTHIC_B, 12, WHITE, True)], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, spacing=1.2)
        if i < n - 1:
            sz = 0.085
            tri = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE,
                                         Inches(bx + bw + (gap - sz) / 2), Inches(y + h / 2 - sz / 2),
                                         Inches(sz), Inches(sz))
            tri.rotation = 90
            tri.shadow.inherit = False
            tri.fill.solid(); tri.fill.fore_color.rgb = GOLD
            tri.line.fill.background()
    if auto: _advance(slide, h)
    return h


def hrule(slide, x, y, w):
    _rect(slide, x, y, w, 0.012, RULE)
